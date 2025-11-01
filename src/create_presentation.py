#!/usr/bin/env python3
"""
Create PowerPoint presentation for Paper 2: Dietary Network Analysis
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def create_presentation():
    """Create comprehensive academic presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
    ACCENT_COLOR = RGBColor(255, 102, 0)  # Orange
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray
    
    def add_title_slide(title, subtitle=""):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullet_points):
        """Add content slide with bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Add line under title
        line = slide.shapes.add_shape(
            1,  # Line shape
            Inches(0.5), Inches(1.1),
            Inches(9), Inches(0)
        )
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)
        
        # Content
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for i, bullet in enumerate(bullet_points):
            if i > 0:
                content_frame.add_paragraph()
            p = content_frame.paragraphs[i]
            
            # Check if it's a sub-bullet (starts with spaces or dash)
            if bullet.startswith('  ') or bullet.startswith('- '):
                p.text = bullet.strip('- ').strip()
                p.level = 1
                p.font.size = Pt(18)
            else:
                p.text = bullet
                p.level = 0
                p.font.size = Pt(22)
            
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(6)
            p.space_after = Pt(6)
        
        return slide
    
    def add_two_column_slide(title, left_content, right_content):
        """Add two-column slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = TITLE_COLOR
        
        # Line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(3)
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4), Inches(5.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, item in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(5.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, item in enumerate(right_content):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        return slide
    
    # Slide 1: Title
    add_title_slide(
        "Personalized Nutrition Through Dietary Network Analysis",
        "Heterogeneity Across Sex, Age, and Metabolic Health\n\n[Your Name] | [Institution] | [Date]"
    )
    
    # Slide 2: Background
    add_content_slide(
        "Background",
        [
            "Metabolic Syndrome (MetS) affects 25-35% of adults worldwide",
            "Dietary patterns are key modifiable risk factors",
            "Traditional approaches analyze individual foods or nutrients",
            "Network science offers new insights into food co-consumption patterns",
            "Gap: Limited understanding of how dietary networks differ across demographic and clinical subgroups"
        ]
    )
    
    # Slide 3: Research Questions
    add_content_slide(
        "Research Questions",
        [
            "How do dietary network patterns differ across:",
            "  Sex (Male vs. Female)",
            "  Age (Young 19-39, Middle 40-59, Older 60-74)",
            "  Metabolic health (MetS+ vs. MetS-)",
            "",
            "Which foods are universal hubs vs. group-specific?",
            "",
            "How do hub foods transition with age?"
        ]
    )
    
    # Slide 4: Study Design
    add_content_slide(
        "Study Design",
        [
            "Data: Korea National Health and Nutrition Examination Survey (KNHANES)",
            "Sample: 22,964 Korean adults (age 19-74)",
            "Stratification: 11 groups (Sex × Age × MetS status)",
            "Food Groups: 12 major categories",
            "Method: Co-occurrence network analysis",
            "Metrics: Degree, Betweenness, Closeness centrality"
        ]
    )
    
    # Slide 5: 12 Food Groups
    add_two_column_slide(
        "12 Food Groups Analyzed",
        [
            "Healthy Foods:",
            "• Grain Products",
            "• Protein Foods",
            "• Vegetables",
            "• Dairy Products",
            "• Fruits"
        ],
        [
            "Unhealthy Foods:",
            "• Fried Foods",
            "• High Fat Meat",
            "• Processed Foods",
            "• Sugar-Sweetened Beverages",
            "• Additional Salt Use",
            "• Salty Food Consumption",
            "• Sweet Food Consumption"
        ]
    )
    
    # Slide 6: Network Construction Method
    add_content_slide(
        "Co-occurrence Network Construction",
        [
            "Step 1: Binary classification (high vs. low consumption)",
            "  Score ≥3 = high, <3 = low",
            "",
            "Step 2: Calculate co-occurrence matrix",
            "  Proportion consuming both foods simultaneously",
            "",
            "Step 3: Apply 70th percentile threshold",
            "  Retain strongest co-occurrence relationships",
            "",
            "Step 4: Create weighted undirected network",
            "  12 nodes (foods), ~20 edges per group"
        ]
    )
    
    # Slide 7: Sample Characteristics
    add_content_slide(
        "Sample Characteristics",
        [
            "Total N = 22,964 adults",
            "  53.5% Male, 46.5% Female",
            "  Mean age: 48.6 ± 11.3 years",
            "",
            "MetS Prevalence: 25.5%",
            "",
            "Group sizes: 516 to 5,629 participants",
            "  Largest: Female middle-aged MetS- (n=5,629)",
            "  Smallest: Male young MetS+ (n=516)"
        ]
    )
    
    # Slide 8: Key Finding 1 - Network Structure
    add_content_slide(
        "Key Finding 1: Consistent Network Structure",
        [
            "All 11 networks have identical topology:",
            "  • 12 nodes (all food groups present)",
            "  • 20 edges (constant)",
            "  • Density = 0.303 (constant)",
            "  • Diameter = 3 (fully connected)",
            "",
            "BUT: Centrality patterns vary substantially",
            "",
            "→ Same structure, different food importance"
        ]
    )
    
    # Slide 9: Key Finding 2 - Universal Hubs
    add_content_slide(
        "Key Finding 2: Three Universal Hub Foods",
        [
            "Present as hubs in ALL 11 groups:",
            "",
            "1. Protein Foods (Top hub, degree: 0.636-1.000)",
            "  Most central in all dietary patterns",
            "",
            "2. Vegetables (Top 3, degree: 0.455-1.000)",
            "  Consistently high across all groups",
            "",
            "3. Grain Products (Top 5, degree: 0.364-0.545)",
            "  Staple food with age-related importance",
            "",
            "→ Core of dietary patterns across all groups"
        ]
    )
    
    # Slide 10: Key Finding 3 - Age-Specific Patterns
    add_content_slide(
        "Key Finding 3: Age-Specific Hub Patterns",
        [
            "Sugar-Sweetened Beverages:",
            "  Young adults: HIGH centrality (0.273-0.364)",
            "  Middle-aged: MODERATE (0.182-0.273)",
            "  Older adults: LOW (0.091-0.182)",
            "  → Dramatic decline with age",
            "",
            "Grain Products (opposite pattern):",
            "  Young adults: Moderate centrality",
            "  Older adults: Highest centrality",
            "  → Progressive increase with age"
        ]
    )
    
    # Slide 11: Key Finding 4 - Sex Differences
    add_content_slide(
        "Key Finding 4: Sex-Specific Patterns",
        [
            "Females:",
            "  • Higher vegetable centrality",
            "  • Sweet food consumption (esp. young)",
            "  • More balanced healthy food patterns",
            "",
            "Males:",
            "  • Higher processed food centrality (MetS+)",
            "  • More fried foods connections (MetS+)",
            "  • Less diverse vegetable consumption"
        ]
    )
    
    # Slide 12: Key Finding 5 - MetS Patterns
    add_content_slide(
        "Key Finding 5: MetS-Specific Patterns",
        [
            "MetS(+) Groups:",
            "  • More connections with unhealthy foods",
            "  • Fried Foods ↔ High Fat Meat",
            "  • Lower fruit & dairy centrality",
            "",
            "MetS(-) Groups:",
            "  • More connections with healthy foods",
            "  • Vegetables ↔ Fruits",
            "  • Dairy ↔ Grain Products",
            "",
            "→ Network structure reflects metabolic health"
        ]
    )
    
    # Slide 13: Hub Transitions
    add_content_slide(
        "Hub Transitions Across Age",
        [
            "Male MetS(+): Young → Middle → Older",
            "  Protein → Protein → Protein (stable)",
            "  Sugary Beverages → Grains → Grains",
            "",
            "Female MetS(-): Young → Middle → Older",
            "  Protein → Protein → Protein",
            "  Sweet Foods → Grains → Grains",
            "",
            "→ Age-related dietary maturation patterns"
        ]
    )
    
    # Slide 14: Clinical Implications - Universal
    add_content_slide(
        "Clinical Implications: Universal Targets",
        [
            "Protein-Vegetable-Grain Triad:",
            "  • Core of all dietary patterns",
            "  • Population-wide intervention target",
            "  • Build meals around this combination",
            "",
            "Fruits as Secondary Target:",
            "  • Promote with vegetables (co-occurrence)",
            "  • Consistently in top 5 hubs",
            "",
            "→ Evidence for population-level guidelines"
        ]
    )
    
    # Slide 15: Clinical Implications - Age-Specific
    add_content_slide(
        "Clinical Implications: Age-Specific",
        [
            "Young Adults (19-39):",
            "  → Reduce sugar-sweetened beverages (high centrality)",
            "  → Replace with healthier alternatives",
            "",
            "Middle-Aged (40-59):",
            "  → Maintain balance during transition",
            "  → Prevent MetS development",
            "",
            "Older Adults (60-74):",
            "  → Leverage grain-centered patterns",
            "  → Add vegetables/fruits to existing meals"
        ]
    )
    
    # Slide 16: Clinical Implications - Personalized
    add_content_slide(
        "Clinical Implications: Personalized Nutrition",
        [
            "Males (especially MetS+):",
            "  → Address processed & fried food connections",
            "  → Practical healthy quick-prep alternatives",
            "",
            "Females:",
            "  → Leverage natural vegetable preference",
            "  → Address sweet food centrality in young",
            "",
            "MetS(+) Individuals:",
            "  → Break unhealthy food co-occurrences",
            "  → Shift toward MetS(-) network patterns"
        ]
    )
    
    # Slide 17: Strengths
    add_content_slide(
        "Study Strengths",
        [
            "Large nationally representative sample (N=22,964)",
            "",
            "Stratified approach revealing heterogeneity",
            "",
            "Co-occurrence networks: interpretable & robust",
            "",
            "Multiple centrality measures (degree, betweenness, closeness)",
            "",
            "Robust across sensitivity analyses",
            "",
            "Reproducible: code and data available"
        ]
    )
    
    # Slide 18: Limitations
    add_content_slide(
        "Limitations",
        [
            "Cross-sectional design → cannot establish causality",
            "",
            "Self-reported dietary data → recall bias",
            "",
            "Food group aggregation → some detail loss",
            "",
            "Binary classification → intensity information lost",
            "",
            "Korean population → generalizability unknown",
            "",
            "Co-occurrence method → other approaches may reveal different insights"
        ]
    )
    
    # Slide 19: Future Directions
    add_content_slide(
        "Future Research Directions",
        [
            "Longitudinal network analysis:",
            "  Track changes over time, causality",
            "",
            "Network-targeted interventions:",
            "  RCTs testing hub food substitution",
            "",
            "Cross-cultural comparisons:",
            "  Western populations, developing countries",
            "",
            "Integration with other data:",
            "  Metabolomics, genomics, microbiome"
        ]
    )
    
    # Slide 20: Conclusions
    add_content_slide(
        "Conclusions",
        [
            "Dietary networks show substantial heterogeneity across sex, age, and MetS status",
            "",
            "Three universal hubs: Protein, Vegetables, Grains",
            "",
            "Group-specific patterns enable personalized nutrition:",
            "  • Age: Sugar drinks → Grains transition",
            "  • Sex: Vegetables (F) vs. Processed foods (M)",
            "  • MetS: Healthy vs. Unhealthy co-occurrences",
            "",
            "Network approach provides actionable insights for both population-wide and individualized interventions"
        ]
    )
    
    # Slide 21: Take-Home Messages
    add_content_slide(
        "Take-Home Messages",
        [
            "🔷 Same food ≠ same importance",
            "   Position in network varies by group",
            "",
            "🔷 Universal + Personalized approach",
            "   Core triad + group-specific tailoring",
            "",
            "🔷 Network thinking in nutrition",
            "   Foods don't exist in isolation",
            "",
            "🔷 Actionable for clinical practice",
            "   Leverage hub foods for intervention"
        ]
    )
    
    # Slide 22: Acknowledgments
    add_content_slide(
        "Acknowledgments & Data Availability",
        [
            "Data Source:",
            "  Korea National Health and Nutrition Examination Survey (KNHANES)",
            "  https://knhanes.kdca.go.kr",
            "",
            "Code & Materials:",
            "  GitHub: https://github.com/HeejeongH/Network",
            "  Network files (GEXF format) available",
            "",
            "Contact:",
            "  [Your Email]",
            "  [Your Institution]"
        ]
    )
    
    # Slide 23: Thank You
    add_title_slide(
        "Thank You!",
        "Questions & Discussion"
    )
    
    # Save presentation
    output_path = '/home/user/webapp/result/Paper2_Dietary_Network_Presentation.pptx'
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    output_file = create_presentation()
    print(f"\n🎉 Presentation created successfully!")
    print(f"📁 Location: {output_file}")
    print(f"\n💡 Recommended use:")
    print(f"   • Academic conferences")
    print(f"   • Research seminars")
    print(f"   • Thesis defense")
    print(f"   • Grant presentations")
